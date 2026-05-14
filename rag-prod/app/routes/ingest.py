import io
import zipfile
import tarfile
import csv
import json
import xml.etree.ElementTree as ET

from fastapi import APIRouter, HTTPException, UploadFile, File
from app.models import IngestRequest, IngestResponse, FileIngestResponse
from app.services.retriever import hybrid_retriever

# Optional imports — graceful fallback
try:
    import fitz
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

try:
    import markdown
    HAS_MARKDOWN = True
except ImportError:
    HAS_MARKDOWN = False

try:
    import ebooklib
    from ebooklib import epub
    HAS_EPUB = True
except ImportError:
    HAS_EPUB = False

router = APIRouter()
CHUNK_SIZE = 500


def extract_pdf(content):
    if not HAS_FITZ:
        raise HTTPException(500, "PyMuPDF not installed. Run: pip install pymupdf")
    pdf = fitz.open(stream=content, filetype="pdf")
    return "\n".join(page.get_text() for page in pdf)


def extract_docx(content):
    if not HAS_DOCX:
        raise HTTPException(500, "python-docx not installed. Run: pip install python-docx")
    doc = DocxDocument(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def extract_pptx(content):
    if not HAS_PPTX:
        raise HTTPException(500, "python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if texts:
            parts.append(f"[Slide {i}] " + " ".join(texts))
    return "\n".join(parts)


def extract_xlsx(content):
    if not HAS_OPENPYXL:
        raise HTTPException(500, "openpyxl not installed. Run: pip install openpyxl")
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[Sheet: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def extract_csv(content):
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("latin-1")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(" | ".join(row) for row in reader if any(c.strip() for c in row))


def extract_json(content):
    try:
        data = json.loads(content.decode("utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        raise HTTPException(400, f"Invalid JSON: {e}")


def extract_xml(content):
    try:
        root = ET.fromstring(content)
        parts = []
        def walk(node, depth=0):
            text = (node.text or "").strip()
            tag = node.tag.split("}")[-1] if "}" in node.tag else node.tag
            if text:
                parts.append("  " * depth + f"<{tag}> {text}")
            for child in node:
                walk(child, depth + 1)
        walk(root)
        return "\n".join(parts)
    except Exception as e:
        raise HTTPException(400, f"Invalid XML: {e}")


def extract_html(content):
    if HAS_BS4:
        soup = BeautifulSoup(content, "html.parser")
        for tag in soup(["script", "style", "meta", "link"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    import re
    text = content.decode("utf-8", errors="ignore")
    return " ".join(re.sub(r"<[^>]+>", " ", text).split())


def extract_markdown(content):
    text = content.decode("utf-8", errors="ignore")
    if HAS_MARKDOWN and HAS_BS4:
        html = markdown.markdown(text)
        return BeautifulSoup(html, "html.parser").get_text(separator="\n", strip=True)
    return text


def extract_plain_text(content):
    for enc in ["utf-8", "utf-16", "latin-1", "cp1252"]:
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def extract_code(content, filename):
    return f"# File: {filename}\n\n{extract_plain_text(content)}"


def extract_epub(content):
    if not HAS_EPUB:
        raise HTTPException(500, "ebooklib not installed. Run: pip install ebooklib")
    book = epub.read_epub(io.BytesIO(content))
    parts = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            if HAS_BS4:
                parts.append(BeautifulSoup(item.get_content(), "html.parser").get_text(separator="\n", strip=True))
            else:
                parts.append(item.get_content().decode("utf-8", errors="ignore"))
    return "\n\n".join(parts)


def extract_zip(content, parent_filename):
    parts = []
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                try:
                    inner = zf.read(name)
                    inner_text = dispatch_extractor(inner, name)
                    if inner_text.strip():
                        parts.append(f"[{name}]\n{inner_text}")
                except Exception as e:
                    parts.append(f"[{name}] (skip: {e})")
    except Exception as e:
        raise HTTPException(400, f"Cannot open ZIP: {e}")
    return "\n\n---\n\n".join(parts)


def extract_tar(content):
    parts = []
    try:
        with tarfile.open(fileobj=io.BytesIO(content)) as tf:
            for member in tf.getmembers():
                if not member.isfile():
                    continue
                try:
                    f = tf.extractfile(member)
                    if f:
                        inner_text = dispatch_extractor(f.read(), member.name)
                        if inner_text.strip():
                            parts.append(f"[{member.name}]\n{inner_text}")
                except Exception as e:
                    parts.append(f"[{member.name}] (skip: {e})")
    except Exception as e:
        raise HTTPException(400, f"Cannot open TAR: {e}")
    return "\n\n---\n\n".join(parts)


CODE_EXTENSIONS = {
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".cpp", ".cc",
    ".cxx", ".h", ".hpp", ".cs", ".go", ".rs", ".rb", ".php", ".swift",
    ".kt", ".kts", ".scala", ".r", ".m", ".sh", ".bash", ".zsh", ".ps1",
    ".bat", ".cmd", ".sql", ".graphql", ".gql", ".proto", ".tf", ".hcl",
    ".dockerfile", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".env", ".vue", ".svelte", ".astro", ".dart", ".lua", ".ex", ".exs",
    ".clj", ".hs", ".ml", ".fsx", ".fs", ".elm", ".nim", ".zig", ".v",
    ".ipynb",
}
TEXT_EXTENSIONS = {".txt", ".log", ".rtf", ".text", ".nfo", ".readme", ".rst"}


def dispatch_extractor(content: bytes, filename: str) -> str:
    name_lower = filename.lower()

    if name_lower.endswith((".tar.gz", ".tgz", ".tar.bz2", ".tar.xz")):
        return extract_tar(content)

    dot = name_lower.rfind(".")
    ext = name_lower[dot:] if dot != -1 else ""

    if ext == ".pdf":       return extract_pdf(content)
    if ext == ".docx":      return extract_docx(content)
    if ext == ".doc":       return extract_plain_text(content)
    if ext == ".pptx":      return extract_pptx(content)
    if ext in (".xlsx", ".xls"): return extract_xlsx(content)
    if ext == ".csv":       return extract_csv(content)
    if ext == ".json":      return extract_json(content)
    if ext == ".xml":       return extract_xml(content)
    if ext in (".html", ".htm"): return extract_html(content)
    if ext in (".md", ".mdx", ".markdown"): return extract_markdown(content)
    if ext == ".epub":      return extract_epub(content)
    if ext == ".zip":       return extract_zip(content, filename)
    if ext == ".tar":       return extract_tar(content)
    if ext in CODE_EXTENSIONS: return extract_code(content, filename)
    if ext in TEXT_EXTENSIONS or ext == "": return extract_plain_text(content)

    try:
        return extract_plain_text(content)
    except Exception:
        raise HTTPException(415, f"Unsupported file type: '{ext}'")


def smart_chunk(text: str, chunk_words: int = CHUNK_SIZE) -> list:
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks, current, current_count = [], [], 0

    for para in paragraphs:
        words = para.split()
        if current_count + len(words) > chunk_words and current:
            chunks.append(" ".join(current))
            current, current_count = [], 0
        current.extend(words)
        current_count += len(words)
        while current_count > chunk_words:
            chunks.append(" ".join(current[:chunk_words]))
            current = current[chunk_words:]
            current_count = len(current)

    if current:
        chunks.append(" ".join(current))
    return [c for c in chunks if c.strip()]


# ─── Routes ──────────────────────────────────────────────────────────────────

@router.post("/ingest", response_model=IngestResponse)
async def ingest_texts(req: IngestRequest):
    try:
        metas = req.metadata or [{} for _ in req.texts]
        hybrid_retriever.add_documents(req.texts, metas)
        return IngestResponse(message="Documents ingested successfully", chunks_stored=len(req.texts))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/ingest/file", response_model=FileIngestResponse)
async def ingest_file(file: UploadFile = File(...)):
    """
    Supports: PDF, DOCX, DOC, PPTX, XLSX, XLS, CSV, JSON, XML,
    HTML, MD, EPUB, TXT, LOG, RTF, RST + all code/config files
    (PY, JS, TS, JSX, TSX, JAVA, C, CPP, GO, RS, RB, PHP, SWIFT,
    KT, CS, SQL, YAML, TOML, SH, PS1, IPYNB, ...) +
    Archives: ZIP, TAR, TAR.GZ (extracts inner files recursively)
    """
    try:
        content = await file.read()
        filename = file.filename or "unknown"
        text = dispatch_extractor(content, filename)

        if not text.strip():
            raise HTTPException(422, f"No text found in '{filename}'.")

        chunks = smart_chunk(text, CHUNK_SIZE)
        if not chunks:
            raise HTTPException(422, "No chunks could be created.")

        metas = [{"source": filename, "chunk": i, "total_chunks": len(chunks)} for i in range(len(chunks))]
        hybrid_retriever.add_documents(chunks, metas)

        return FileIngestResponse(message="File ingested successfully", chunks_stored=len(chunks), filename=filename)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))